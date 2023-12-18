""" 
Ref for slide types: 
0 ->  title and subtitle
1 ->  title and content
        2 ->  section header
        3 ->  two content
        4 ->  Comparison
        5 ->  Title only 
        6 ->  Blank
        7 ->  Content with caption
        8 ->  Pic with caption
"""

from pptx import Presentation
import json
import os

def create_pptx(textInput,filename,isGPT):
    if not textInput:
        return False
    if not filename:
        return False
    if(isGPT==False):
        slides = textInput.split('\n\n')
    else:
        slides = json.loads(textInput)
    
    return createPresentation(slides,filename,isGPT)

def create_presentation(slides,filename,isGPT):
    print(slides,type(slides))
    prs = Presentation()
    
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Slides"
    subtitle.text = "Created by Presentai"

    for slide in slides:
        new_slide=prs.slides.add_slide(prs.slide_layouts[1])
        #new_slide.shapes.title.text=slide['title']
        tf=new_slide.shapes.placeholders[1].text_frame
        if(isGPT==False):
            items = slide.split('.')
            for item in items:
                tf.add_paragraph().text=item
        else:
            new_slide.shapes.title.text=slide['title']        
            items = slide['content'].split('.')
            for item in items:
                tf.add_paragraph().text=item
    filepath=os.path.join('userDocs',filename)     
    prs.save(filepath)
    return True